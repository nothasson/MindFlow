import io
import re
from html.parser import HTMLParser
from urllib import request

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation


class HTMLContentExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.title = ""
        self._texts: list[str] = []
        self._in_title = False
        self._skip_depth = 0

    def handle_starttag(self, tag, attrs):
        normalized = tag.lower()
        if normalized in {"script", "style", "noscript"}:
            self._skip_depth += 1
        if normalized == "title":
            self._in_title = True

    def handle_endtag(self, tag):
        normalized = tag.lower()
        if normalized in {"script", "style", "noscript"} and self._skip_depth > 0:
            self._skip_depth -= 1
        if normalized == "title":
            self._in_title = False

    def handle_data(self, data):
        text = normalize_whitespace(data)
        if not text:
            return
        if self._in_title and not self.title:
            self.title = text
            return
        if self._skip_depth == 0:
            self._texts.append(text)

    @property
    def text(self) -> str:
        return normalize_whitespace("\n\n".join(self._texts))


def parse_pdf(file_bytes: bytes, filename: str) -> dict:
    """解析 PDF 文件，提取纯文本"""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()

    full_text = "\n\n".join(pages)
    return {
        "text": full_text,
        "pages": len(pages),
        "filename": filename,
    }


def parse_docx(file_bytes: bytes, filename: str) -> dict:
    """解析 Word (.docx) 文档，提取段落文本"""
    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)
    return {
        "text": full_text,
        "pages": max(len(paragraphs), 1),
        "filename": filename,
    }


def parse_pptx(file_bytes: bytes, filename: str) -> dict:
    """解析 PPT (.pptx) 文件，提取每页幻灯片的文本"""
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text: list[str] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append("\n".join(texts))
    full_text = "\n\n".join(slides_text)
    return {
        "text": full_text,
        "pages": len(prs.slides),
        "filename": filename,
    }


def parse_text(content: str, filename: str) -> dict:
    """处理纯文本文件"""
    return {
        "text": content,
        "pages": 1,
        "filename": filename,
    }


def parse_url(url: str) -> dict:
    """抓取并解析网页正文。"""
    html = fetch_url_content(url)
    extractor = HTMLContentExtractor()
    extractor.feed(html)
    extractor.close()

    title = extractor.title or url
    text = extractor.text
    return {
        "text": text,
        "title": title,
        "source_url": url,
    }


def fetch_url_content(url: str) -> str:
    req = request.Request(
        url,
        headers={
            "User-Agent": "MindFlow/1.0 (+https://github.com/nothasson/MindFlow)",
            "Accept": "text/html,application/xhtml+xml",
        },
    )
    with request.urlopen(req, timeout=30) as resp:
        charset = resp.headers.get_content_charset() or "utf-8"
        return resp.read().decode(charset, errors="replace")


def normalize_whitespace(text: str) -> str:
    text = text.replace("\xa0", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()
