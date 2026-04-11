"""parser 模块单元测试：覆盖 DOCX / PPTX 解析。"""

import io

from docx import Document as DocxDocument
from pptx import Presentation

from app.services.parser import parse_docx, parse_pptx, parse_text


def _make_docx(paragraphs: list[str]) -> bytes:
    """创建最小 DOCX 文件并返回字节内容。"""
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(slides_texts: list[list[str]]) -> bytes:
    """创建最小 PPTX 文件并返回字节内容。
    slides_texts: 每个元素是一页幻灯片里的文本列表。
    """
    prs = Presentation()
    for texts in slides_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        for i, text in enumerate(texts):
            from pptx.util import Inches
            left = Inches(1)
            top = Inches(1 + i)
            txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(1))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------- DOCX 测试 ----------


def test_parse_docx_basic():
    """解析包含多个段落的 DOCX 文件。"""
    content = _make_docx(["第一段内容", "第二段内容", "第三段内容"])
    result = parse_docx(content, "test.docx")

    assert result["filename"] == "test.docx"
    assert "第一段内容" in result["text"]
    assert "第二段内容" in result["text"]
    assert "第三段内容" in result["text"]
    assert result["pages"] == 3  # 3 个非空段落


def test_parse_docx_empty():
    """解析空 DOCX 文件应返回空文本，pages 至少为 1。"""
    content = _make_docx([])
    result = parse_docx(content, "empty.docx")

    assert result["text"] == ""
    assert result["pages"] >= 1
    assert result["filename"] == "empty.docx"


def test_parse_docx_with_blank_paragraphs():
    """包含空白段落的 DOCX，空白段落应被过滤。"""
    content = _make_docx(["有内容", "", "  ", "另一段"])
    result = parse_docx(content, "blanks.docx")

    assert "有内容" in result["text"]
    assert "另一段" in result["text"]
    # pages 只计非空段落
    assert result["pages"] == 2


# ---------- PPTX 测试 ----------


def test_parse_pptx_basic():
    """解析包含多页幻灯片的 PPTX 文件。"""
    content = _make_pptx([
        ["标题一", "正文内容 A"],
        ["标题二", "正文内容 B"],
    ])
    result = parse_pptx(content, "test.pptx")

    assert result["filename"] == "test.pptx"
    assert "标题一" in result["text"]
    assert "正文内容 A" in result["text"]
    assert "标题二" in result["text"]
    assert result["pages"] == 2


def test_parse_pptx_empty_slides():
    """空 PPTX（无幻灯片）应返回空文本。"""
    prs = Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    content = buf.getvalue()

    result = parse_pptx(content, "empty.pptx")

    assert result["text"] == ""
    assert result["pages"] == 0
    assert result["filename"] == "empty.pptx"


# ---------- 纯文本测试 ----------


def test_parse_text():
    """解析纯文本应原样返回。"""
    result = parse_text("这是一段学习资料。", "note.txt")

    assert result["text"] == "这是一段学习资料。"
    assert result["pages"] == 1
    assert result["filename"] == "note.txt"
