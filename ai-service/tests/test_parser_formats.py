"""测试多格式文档解析：DOCX 和 PPTX"""

import io

from docx import Document as DocxDocument
from pptx import Presentation

from app.services.parser import parse_docx, parse_pptx


def _make_docx(paragraphs: list[str]) -> bytes:
    """创建最小 DOCX 文件并返回字节"""
    doc = DocxDocument()
    for text in paragraphs:
        doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(slides_texts: list[list[str]]) -> bytes:
    """创建最小 PPTX 文件并返回字节。
    slides_texts: 每个元素是一页幻灯片的文本列表。
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # 空白布局
    for texts in slides_texts:
        slide = prs.slides.add_slide(blank_layout)
        for i, text in enumerate(texts):
            from pptx.util import Inches
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(1))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_docx_extracts_paragraphs():
    """验证 DOCX 解析能正确提取段落文本"""
    paragraphs = ["第一章 线性代数基础", "向量空间是线性代数的核心概念。", "矩阵运算包括加法和乘法。"]
    content = _make_docx(paragraphs)

    result = parse_docx(content, "test.docx")

    assert result["filename"] == "test.docx"
    assert "线性代数基础" in result["text"]
    assert "向量空间" in result["text"]
    assert "矩阵运算" in result["text"]
    assert result["pages"] >= 1


def test_parse_docx_empty():
    """验证空 DOCX 返回空文本"""
    content = _make_docx([])

    result = parse_docx(content, "empty.docx")

    assert result["filename"] == "empty.docx"
    assert result["text"] == ""
    assert result["pages"] >= 1


def test_parse_pptx_extracts_slides():
    """验证 PPTX 解析能正确提取幻灯片文本"""
    slides = [
        ["第一页标题", "第一页内容"],
        ["第二页标题"],
    ]
    content = _make_pptx(slides)

    result = parse_pptx(content, "test.pptx")

    assert result["filename"] == "test.pptx"
    assert result["pages"] == 2
    assert "第一页标题" in result["text"]
    assert "第一页内容" in result["text"]
    assert "第二页标题" in result["text"]


def test_parse_pptx_empty():
    """验证空 PPTX 返回空文本"""
    prs = Presentation()
    buf = io.BytesIO()
    prs.save(buf)
    content = buf.getvalue()

    result = parse_pptx(content, "empty.pptx")

    assert result["filename"] == "empty.pptx"
    assert result["text"] == ""
    assert result["pages"] == 0
